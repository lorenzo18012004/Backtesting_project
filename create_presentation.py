"""
Génère une présentation PowerPoint du projet Backtesting.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Couleurs
ACCENT = RGBColor(59, 130, 246)     # Bleu
TEXT_DARK = RGBColor(15, 23, 42)
TEXT_MUTED = RGBColor(100, 116, 139)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = TEXT_MUTED
        p2.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(title, bullets, sub_bullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    # Content
    tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.5))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        para = tf2.paragraphs[i] if i < len(tf2.paragraphs) else tf2.add_paragraph()
        para.text = f"• {bullet}"
        para.font.size = Pt(18)
        para.font.color.rgb = TEXT_DARK
        para.space_before = Pt(12)
        if sub_bullets and i < len(sub_bullets) and sub_bullets[i]:
            for sub in sub_bullets[i]:
                p2 = tf2.add_paragraph()
                p2.text = f"    → {sub}"
                p2.font.size = Pt(14)
                p2.font.color.rgb = TEXT_MUTED
                p2.space_before = Pt(4)
    return slide

def add_code_slide(title, code):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = code
    p2.font.size = Pt(14)
    p2.font.name = "Consolas"
    p2.font.color.rgb = TEXT_DARK
    return slide

# Slide 1 - Titre
add_title_slide("Backtesting Engine", "Moteur de backtesting pour stratégies de trading — Lorenzo PHILIPPE")

# Slide 2 - Le piège du backtesting
add_content_slide(
    "Le piège du backtesting",
    [
        "Une stratégie qui « marche » en backtest peut échouer en réel.",
        "Pourquoi ?",
    ],
    [
        ["Overfitting : trop de paramètres, la stratégie s’adapte au bruit"],
        ["Look-ahead : utilisation d’informations futures sans le savoir"],
        ["Frais et slippage oubliés : les coûts réduisent les performances"],
        ["Période courte ou unique : pas de validation sur d’autres marchés"],
    ]
)

# Slide 3 - La solution
add_content_slide(
    "Ce que j’ai mis en place",
    [
        "Un moteur qui intègre les garde-fous pour limiter ces biais.",
        "Pipeline en 7 étapes : Data → Nettoyage → Signaux → PnL → Frais → Risque → Viz",
    ],
    [
        [],
        ["Shift anti look-ahead", "Frais + slippage par trade", "Walk-forward (in/out-of-sample)", "Monte Carlo & Bootstrap pour la robustesse"],
    ]
)

# Slide 4 - Stratégies
add_content_slide(
    "Les 3 stratégies + Live",
    [
        "Buy & Hold — Référence, toujours long",
        "SMA Crossover + RSI — Croisement de moyennes mobiles avec filtre RSI",
        "Portefeuille multi-actifs — Style hedge fund : Markowitz, facteurs (momentum, trend, low vol), VaR, circuit breaker",
        "Live — Suivi en temps réel de la stratégie portefeuille",
    ],
)

# Slide 5 - Métriques de risque
add_content_slide(
    "Métriques de risque",
    [
        "Sharpe, Sortino, Calmar — Rendement ajusté au risque",
        "VaR, Expected Shortfall — Risque de queue",
        "Max Drawdown, Win Rate, Profit Factor",
        "Monte Carlo, Bootstrap — Robustesse des résultats",
        "Walk-forward — Validation out-of-sample",
    ],
)

# Slide 6 - Architecture
add_content_slide(
    "Architecture du code",
    [
        "backtest/ — Package modulaire",
    ],
    [
        [
            "config.py — Constantes",
            "data.py — Fetch Yahoo Finance, nettoyage, log-returns",
            "signals.py — Génération de signaux (SMA, RSI, etc.)",
            "pnl.py — Rendements, frais, slippage",
            "risk.py — Sharpe, VaR, Monte Carlo…",
            "portfolio.py — Stratégie multi-actifs",
            "viz.py — Graphiques",
            "core.py — Orchestration",
        ],
    ]
)

# Slide 7 - Comment lancer
add_code_slide(
    "Comment lancer l’app",
    """git clone https://github.com/lorenzo18012004/Backtesting_project.git
cd Backtesting_project
pip install -r requirements.txt
streamlit run app.py

→ L'app s'ouvre sur http://localhost:8501"""
)

# Slide 8 - Conclusion
add_content_slide(
    "Conclusion",
    [
        "Objectif : comprendre pourquoi une stratégie peut sembler gagnante en backtest mais échouer en réel.",
        "Ce projet intègre les garde-fous pour limiter ces biais.",
        "Code open source — Retours et suggestions bienvenus.",
    ],
)

# Slide 9 - Liens
add_content_slide(
    "Liens",
    [
        "GitHub : https://github.com/lorenzo18012004/Backtesting_project",
        "Stack : Python, Pandas, NumPy, SciPy, Streamlit, yfinance",
        "29 tests unitaires et d’intégration",
    ],
)

prs.save("Backtesting_Engine_Presentation.pptx")
print("Présentation créée : Backtesting_Engine_Presentation.pptx")
